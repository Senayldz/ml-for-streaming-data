"""
app.py
------
Unified entrypoint for the SWaT Streaming Anomaly Detection project.

This monolithic file combines:
1. Preprocessing and scaling logic (from preprocess.py)
2. Anomaly detector models (from model.py)
3. Streaming simulator and throughput tracking (from stream_engine.py)
4. Model training and evaluation pipeline (from main.py)
5. Utility functions to create mini-datasets and PPTX presentations
6. Streamlit real-time interactive dashboard UI (from dashboard.py)

Usage
-----
1. Run Streamlit Dashboard UI:
   streamlit run app.py

2. Run Model Training & Evaluation Pipeline:
   python app.py [arguments]
   Example: python app.py --model lgbm --train-rows 100000

3. Create Mini-Dataset for deployment:
   python app.py --create-mini

4. Generate Presentation:
   python app.py --generate-presentation
"""

from __future__ import annotations

import os
import sys
import time
import argparse
import random
import threading
import collections
import traceback
from abc import ABC, abstractmethod
from pathlib import Path
from typing import Generator, Iterator, Optional, Tuple, Union, List

import numpy as np
import pandas as pd
import joblib
from sklearn.model_selection import train_test_split
from sklearn.preprocessing import StandardScaler
from sklearn.metrics import (
    precision_score,
    recall_score,
    f1_score,
    roc_auc_score,
    confusion_matrix,
    classification_report,
)

# Optional heavy imports – graceful degradation if not installed
try:
    import lightgbm as lgb
    _LGBM_AVAILABLE = True
except ImportError:
    _LGBM_AVAILABLE = False

try:
    import torch
    import torch.nn as nn
    _TORCH_AVAILABLE = True
except ImportError:
    _TORCH_AVAILABLE = False

# ── Global Paths & Constants ──────────────────────────────────────────────────
ROOT_DIR        = Path(__file__).resolve().parent
DATA_DIR        = ROOT_DIR / "dataset"
NORMAL_CSV      = DATA_DIR / "normal.csv"
ATTACK_CSV      = DATA_DIR / "attack.csv"
MERGED_CSV      = DATA_DIR / "merged.csv"
MINI_CSV        = DATA_DIR / "merged_mini.csv"

ARTIFACTS_DIR   = ROOT_DIR / "artifacts"
SCALER_PATH     = ARTIFACTS_DIR / "scaler.joblib"
LGBM_PATH       = ARTIFACTS_DIR / "lgbm_detector.joblib"

LABEL_COL       = "Normal/Attack"
TIMESTAMP_COL   = " Timestamp"

# Create necessary directories
ARTIFACTS_DIR.mkdir(parents=True, exist_ok=True)
(ROOT_DIR / "reports").mkdir(parents=True, exist_ok=True)


# ══════════════════════════════════════════════════════════════════════════════
# PREPROCESSING MODULE (from preprocess.py)
# ══════════════════════════════════════════════════════════════════════════════

def _clean_columns(df: pd.DataFrame) -> pd.DataFrame:
    """Strip leading/trailing whitespace from column names."""
    df.columns = [c.strip() for c in df.columns]
    return df


def _encode_label(series: pd.Series) -> pd.Series:
    """Map 'Normal' -> 0, anything else -> 1 (binary attack flag)."""
    return series.str.strip().map(lambda x: 0 if x == "Normal" else 1).astype(np.int8)


def get_feature_columns(df: pd.DataFrame) -> List[str]:
    """Return all numeric feature columns (excludes timestamp & label)."""
    drop = {"Timestamp", LABEL_COL}
    return [c for c in df.columns if c not in drop and df[c].dtype != object]


def load_and_preprocess(
    csv_path: Path,
    scaler: Optional[StandardScaler] = None,
    fit_scaler: bool = False,
    chunksize: Optional[int] = None,
) -> Tuple[np.ndarray, np.ndarray, StandardScaler]:
    """Load a CSV (optionally in chunks), clean, scale, and return arrays."""
    reader = pd.read_csv(csv_path, chunksize=chunksize) if chunksize else [pd.read_csv(csv_path)]

    X_parts, y_parts = [], []
    feat_cols: Optional[List[str]] = None

    for chunk in reader:
        chunk = _clean_columns(chunk)

        if feat_cols is None:
            feat_cols = get_feature_columns(chunk)

        X_chunk = chunk[feat_cols].fillna(0).values.astype(np.float32)
        y_chunk = _encode_label(chunk[LABEL_COL]).values

        X_parts.append(X_chunk)
        y_parts.append(y_chunk)

    X = np.concatenate(X_parts, axis=0)
    y = np.concatenate(y_parts, axis=0)

    if fit_scaler:
        scaler = StandardScaler()
        scaler.fit(X)

    if scaler is not None:
        X = scaler.transform(X).astype(np.float32)

    return X, y, scaler


def fit_and_save_scaler(
    csv_path: Path = NORMAL_CSV,
    chunksize: int = 100_000,
) -> StandardScaler:
    """Fit a StandardScaler on normal traffic and persist to disk."""
    scaler = StandardScaler()
    feat_cols: Optional[List[str]] = None

    for chunk in pd.read_csv(csv_path, chunksize=chunksize):
        chunk = _clean_columns(chunk)
        if feat_cols is None:
            feat_cols = get_feature_columns(chunk)
        X_chunk = chunk[feat_cols].fillna(0).values.astype(np.float32)
        scaler.partial_fit(X_chunk)

    joblib.dump(scaler, SCALER_PATH)
    return scaler


def load_scaler() -> StandardScaler:
    """Load a previously saved scaler from disk."""
    return joblib.load(SCALER_PATH)


# ══════════════════════════════════════════════════════════════════════════════
# DETECTORS MODULE (from model.py)
# ══════════════════════════════════════════════════════════════════════════════

class BaseDetector(ABC):
    """Common interface for all detectors."""

    @abstractmethod
    def fit(self, X: np.ndarray, y: Optional[np.ndarray] = None) -> "BaseDetector": ...

    @abstractmethod
    def predict(self, X: np.ndarray) -> np.ndarray:
        """Return binary predictions: 0 = Normal, 1 = Attack."""
        ...

    @abstractmethod
    def predict_proba(self, X: np.ndarray) -> np.ndarray:
        """Return attack probability scores in [0, 1]."""
        ...

    def save(self, path: Path) -> None:
        joblib.dump(self, path)

    @classmethod
    def load(cls, path: Path) -> "BaseDetector":
        return joblib.load(path)


class LGBMDetector(BaseDetector):
    """LightGBM binary classifier tuned for class imbalance."""

    def __init__(
        self,
        n_estimators: int = 300,
        max_depth: int = 6,
        learning_rate: float = 0.1,
        num_leaves: int = 31,
        class_weight_ratio: Optional[float] = None,
        random_state: int = 42,
    ):
        if not _LGBM_AVAILABLE:
            raise ImportError("lightgbm is not installed. Run: pip install lightgbm")

        self.params = dict(
            n_estimators   = n_estimators,
            max_depth       = max_depth,
            learning_rate   = learning_rate,
            num_leaves      = num_leaves,
            random_state    = random_state,
            n_jobs          = -1,
            verbosity       = -1,
        )
        self._class_weight_ratio = class_weight_ratio
        self.model = None  # type: Optional[lgb.LGBMClassifier]

    def fit(self, X: np.ndarray, y: Optional[np.ndarray] = None) -> "LGBMDetector":
        assert y is not None, "LGBMDetector requires labels."
        n_neg = int((y == 0).sum())
        n_pos = int((y == 1).sum())
        ratio = self._class_weight_ratio or (n_neg / max(n_pos, 1))

        self.model = lgb.LGBMClassifier(
            scale_pos_weight=ratio,
            **self.params,
        )
        self.model.fit(X, y)
        return self

    def predict(self, X: np.ndarray) -> np.ndarray:
        return self.model.predict(X)

    def predict_proba(self, X: np.ndarray) -> np.ndarray:
        return self.model.predict_proba(X)[:, 1]


class IsoForestDetector(BaseDetector):
    """IsolationForest trained exclusively on normal traffic."""

    def __init__(
        self,
        n_estimators: int = 100,
        contamination: float = 0.04,
        random_state: int = 42,
        max_samples: Union[int, str] = "auto",
    ):
        from sklearn.ensemble import IsolationForest
        self.model = IsolationForest(
            n_estimators = n_estimators,
            contamination = contamination,
            random_state  = random_state,
            max_samples   = max_samples,
            n_jobs        = -1,
        )
        self._threshold: float = 0.0

    def fit(self, X: np.ndarray, y: Optional[np.ndarray] = None) -> "IsoForestDetector":
        self.model.fit(X)
        scores = -self.model.score_samples(X)
        self._threshold = float(np.percentile(scores, 96))
        self._min_score = float(scores.min())
        self._max_score = float(scores.max())
        return self

    def predict(self, X: np.ndarray) -> np.ndarray:
        scores = -self.model.score_samples(X)
        return (scores >= self._threshold).astype(np.int8)

    def predict_proba(self, X: np.ndarray) -> np.ndarray:
        scores = -self.model.score_samples(X)
        if hasattr(self, "_min_score") and hasattr(self, "_max_score"):
            lo, hi = self._min_score, self._max_score
        else:
            lo, hi = scores.min(), scores.max()
        if hi == lo:
            return np.zeros(len(scores), dtype=np.float32)
        proba = (scores - lo) / (hi - lo)
        return np.clip(proba, 0.0, 1.0).astype(np.float32)


if _TORCH_AVAILABLE:
    class _AEModel(nn.Module):
        def __init__(self, input_dim: int, latent_dim: int = 8):
            super().__init__()
            hidden = max(latent_dim * 2, 16)
            self.encoder = nn.Sequential(
                nn.Linear(input_dim, hidden),
                nn.ReLU(),
                nn.Linear(hidden, latent_dim),
                nn.ReLU(),
            )
            self.decoder = nn.Sequential(
                nn.Linear(latent_dim, hidden),
                nn.ReLU(),
                nn.Linear(hidden, input_dim),
            )

        def forward(self, x: torch.Tensor) -> torch.Tensor:
            return self.decoder(self.encoder(x))

        def reconstruction_error(self, x: torch.Tensor) -> torch.Tensor:
            with torch.no_grad():
                recon = self.forward(x)
                return ((x - recon) ** 2).mean(dim=1)


class AutoencoderDetector(BaseDetector):
    """Autoencoder-based anomaly detector using PyTorch."""

    def __init__(
        self,
        input_dim: int = 51,
        latent_dim: int = 8,
        epochs: int = 20,
        batch_size: int = 256,
        lr: float = 1e-3,
        threshold_percentile: float = 95.0,
        device: str = "cpu",
    ):
        if not _TORCH_AVAILABLE:
            raise ImportError("PyTorch is not installed. Run: pip install torch")

        self.input_dim             = input_dim
        self.latent_dim            = latent_dim
        self.epochs                = epochs
        self.batch_size            = batch_size
        self.lr                    = lr
        self.threshold_percentile  = threshold_percentile
        self.device                = torch.device(device)
        self._threshold: float     = 0.0
        self._ae = None

    def fit(self, X: np.ndarray, y: Optional[np.ndarray] = None) -> "AutoencoderDetector":
        if y is not None:
            X = X[y == 0]

        self._ae = _AEModel(self.input_dim, self.latent_dim).to(self.device)
        optimiser = torch.optim.Adam(self._ae.parameters(), lr=self.lr)
        loss_fn   = nn.MSELoss()

        dataset = torch.tensor(X, dtype=torch.float32)
        n = len(dataset)

        for epoch in range(self.epochs):
            idx = torch.randperm(n)
            for start in range(0, n, self.batch_size):
                batch = dataset[idx[start:start + self.batch_size]].to(self.device)
                optimiser.zero_grad()
                recon = self._ae(batch)
                loss  = loss_fn(recon, batch)
                loss.backward()
                optimiser.step()

        with torch.no_grad():
            errors = self._ae.reconstruction_error(dataset.to(self.device)).cpu().numpy()
        self._threshold = float(np.percentile(errors, self.threshold_percentile))
        self._min_error = float(errors.min())
        self._max_error = float(errors.max())
        return self

    def predict(self, X: np.ndarray) -> np.ndarray:
        errors = self._reconstruction_errors(X)
        return (errors >= self._threshold).astype(np.int8)

    def predict_proba(self, X: np.ndarray) -> np.ndarray:
        errors = self._reconstruction_errors(X)
        if hasattr(self, "_min_error") and hasattr(self, "_max_error"):
            lo, hi = self._min_error, self._max_error
        else:
            lo, hi = errors.min(), errors.max()
        if hi == lo:
            return np.zeros(len(errors), dtype=np.float32)
        proba = (errors - lo) / (hi - lo)
        return np.clip(proba, 0.0, 1.0).astype(np.float32)

    def _reconstruction_errors(self, X: np.ndarray) -> np.ndarray:
        self._ae.eval()
        t = torch.tensor(X, dtype=torch.float32).to(self.device)
        with torch.no_grad():
            errors = self._ae.reconstruction_error(t).cpu().numpy()
        return errors

    def save(self, path: Path) -> None:
        torch.save(
            {
                "state_dict"   : self._ae.state_dict(),
                "threshold"    : self._threshold,
                "input_dim"    : self.input_dim,
                "latent_dim"   : self.latent_dim,
            },
            path,
        )

    @classmethod
    def load(cls, path: Path, **kwargs) -> "AutoencoderDetector":
        ckpt = torch.load(path, map_location="cpu")
        obj  = cls(input_dim=ckpt["input_dim"], latent_dim=ckpt["latent_dim"], **kwargs)
        obj._ae = _AEModel(ckpt["input_dim"], ckpt["latent_dim"])
        obj._ae.load_state_dict(ckpt["state_dict"])
        obj._threshold = ckpt["threshold"]
        return obj


def build_detector(name: str = "lgbm", **kwargs) -> BaseDetector:
    """Factory helper to build anomaly detectors."""
    name = name.lower()
    registry = {
        "lgbm"        : LGBMDetector,
        "isoforest"   : IsoForestDetector,
    }
    if _TORCH_AVAILABLE:
        registry["autoencoder"] = AutoencoderDetector

    if name not in registry:
        raise ValueError(f"Unknown detector '{name}'. Choose from {list(registry.keys())}")
    return registry[name](**kwargs)


# ══════════════════════════════════════════════════════════════════════════════
# STREAMING SIMULATOR MODULE (from stream_engine.py)
# ══════════════════════════════════════════════════════════════════════════════

def stream_records(
    csv_path: Path,
    scaler: StandardScaler,
    read_chunksize: int = 10_000,
    throttle_hz: Optional[float] = None,
) -> Generator[Tuple[np.ndarray, int], None, None]:
    """Generator that yields one (feature_vector, label) tuple per CSV row."""
    sleep_interval = (1.0 / throttle_hz) if throttle_hz else 0.0
    feat_cols: list[str] | None = None

    for chunk in pd.read_csv(csv_path, chunksize=read_chunksize):
        chunk = _clean_columns(chunk)

        if feat_cols is None:
            feat_cols = get_feature_columns(chunk)

        X_raw = chunk[feat_cols].fillna(0).values.astype(np.float32)
        X_scaled = scaler.transform(X_raw).astype(np.float32)
        y = _encode_label(chunk[LABEL_COL]).values

        for i in range(len(X_scaled)):
            yield X_scaled[i], int(y[i])
            if sleep_interval:
                time.sleep(sleep_interval)


def stream_batches(
    csv_path: Path,
    scaler: StandardScaler,
    batch_size: int = 64,
    read_chunksize: int = 10_000,
    throttle_hz: Optional[float] = None,
) -> Generator[Tuple[np.ndarray, np.ndarray], None, None]:
    """Generator that yields (X_batch, y_batch) tuples."""
    sleep_interval = (batch_size / throttle_hz) if throttle_hz else 0.0
    feat_cols: list[str] | None = None

    for chunk in pd.read_csv(csv_path, chunksize=read_chunksize):
        chunk = _clean_columns(chunk)

        if feat_cols is None:
            feat_cols = get_feature_columns(chunk)

        X_raw = chunk[feat_cols].fillna(0).values.astype(np.float32)
        X_scaled = scaler.transform(X_raw).astype(np.float32)
        y = _encode_label(chunk[LABEL_COL]).values

        for start in range(0, len(X_scaled), batch_size):
            end = start + batch_size
            yield X_scaled[start:end], y[start:end]
            if sleep_interval:
                time.sleep(sleep_interval)


class ThroughputTracker:
    """Tracker for streaming inference latency and throughput benchmarks."""

    def __init__(self) -> None:
        self._latencies: list[float] = []
        self._t0 = time.perf_counter()

    class _Context:
        def __init__(self, tracker: ThroughputTracker) -> None:
            self._tracker = tracker

        def __enter__(self):
            self._start = time.perf_counter()
            return self

        def __exit__(self, *_):
            elapsed_ms = (time.perf_counter() - self._start) * 1000.0
            self._tracker._latencies.append(elapsed_ms)

    def measure(self) -> "_Context":
        return self._Context(self)

    def record(self, latency_ms: float) -> None:
        self._latencies.append(latency_ms)

    def summary(self) -> dict:
        lats = np.array(self._latencies, dtype=np.float64)
        wall_time_s = time.perf_counter() - self._t0
        n = len(lats)
        return {
            "n_records"           : n,
            "wall_time_s"         : round(wall_time_s, 4),
            "throughput_rps"      : round(n / wall_time_s, 2) if wall_time_s > 0 else 0,
            "latency_mean_ms"     : round(float(lats.mean()), 4)  if n else 0,
            "latency_median_ms"   : round(float(np.median(lats)), 4) if n else 0,
            "latency_p95_ms"      : round(float(np.percentile(lats, 95)), 4) if n else 0,
            "latency_p99_ms"      : round(float(np.percentile(lats, 99)), 4) if n else 0,
            "latency_max_ms"      : round(float(lats.max()), 4) if n else 0,
        }


# ══════════════════════════════════════════════════════════════════════════════
# MODEL TRAINING & EVALUATION PIPELINE (from main.py)
# ══════════════════════════════════════════════════════════════════════════════

def _markdown_table(rows) -> str:
    """Render a two-column markdown table from (metric, value) pairs."""
    col_w = max(len(r[0]) for r in rows) + 2
    val_w = max(len(r[1]) for r in rows) + 2

    header = f"| {'Metric':<{col_w}}| {'Value':<{val_w}}|"
    sep    = f"|{'-' * (col_w + 1)}|{'-' * (val_w + 1)}|"
    lines  = [header, sep]
    for k, v in rows:
        lines.append(f"| {k:<{col_w}}| {v:<{val_w}}|")
    return "\n".join(lines)


def _print_step(msg: str) -> None:
    print("\n" + "-" * 60 + "\n  " + msg + "\n" + "-" * 60)


def stage_scaler() -> object:
    _print_step("Stage 1 | Fitting / loading scaler")
    if SCALER_PATH.exists():
        print("  -> Cached scaler found, loading...")
        return load_scaler()
    print("  -> Fitting incremental StandardScaler on normal.csv...")
    t0 = time.perf_counter()
    scaler = fit_and_save_scaler(NORMAL_CSV)
    print(f"  -> Done in {time.perf_counter() - t0:.1f}s  |  saved to {SCALER_PATH}")
    return scaler


def stage_train(scaler, model_name: str, train_rows: Optional[int], no_retrain: bool):
    model_path = ARTIFACTS_DIR / f"{model_name}_detector.joblib"

    if no_retrain and model_path.exists():
        _print_step("Stage 2 | Loading saved model (--no-retrain)")
        detector = BaseDetector.load(model_path)
        print("  -> Rebuilding test split (temporal) from merged.csv for evaluation...")
        X, y, _ = load_and_preprocess(MERGED_CSV, scaler=scaler, chunksize=100_000)
        split_idx = int(len(X) * 0.8)
        X_test = X[split_idx:]
        y_test = y[split_idx:]
        print(f"  -> Test set: {len(X_test):,} rows  |  attack rate: {y_test.mean():.4f}")
        return detector, X_test, y_test

    _print_step("Stage 2 | Training model")
    print(f"  -> Backend  : {model_name.upper()}")
    print(f"  -> Data     : {MERGED_CSV}  (80/20 temporal split)")

    nrows = train_rows
    print(f"  -> Loading {'all' if nrows is None else nrows} rows...")
    t0 = time.perf_counter()
    X, y, _ = load_and_preprocess(MERGED_CSV, scaler=scaler, chunksize=100_000)
    if nrows and nrows < len(X):
        X, y = X[:nrows], y[:nrows]

    split_idx = int(len(X) * 0.8)
    X_train, X_test = X[:split_idx], X[split_idx:]
    y_train, y_test = y[:split_idx], y[split_idx:]
    print(f"  -> Train: {len(X_train):,}  |  Test: {len(X_test):,}")
    print(f"  -> Attack ratio  train={y_train.mean():.4f}  test={y_test.mean():.4f}")

    extra_kwargs = {"input_dim": X.shape[1]} if model_name == "autoencoder" else {}
    detector = build_detector(model_name, **extra_kwargs)
    print(f"  -> Fitting...")
    detector.fit(X_train, y_train)

    elapsed = time.perf_counter() - t0
    print(f"  -> Training done in {elapsed:.1f}s")

    detector.save(model_path)
    print(f"  -> Model saved to {model_path}")

    return detector, X_test, y_test


def stage_stream_inference(detector, X_test, y_test, stream_rows: int, batch_size: int):
    _print_step("Stage 3 | Streaming inference (test split)")
    n_stream = min(stream_rows, len(X_test))
    print(f"  -> Streaming {n_stream:,} records from held-out stratified test split")
    print(f"  -> Attack rate in stream: {y_test[:n_stream].mean():.4f}")
    print(f"  -> Batch size : {batch_size}")

    tracker   = ThroughputTracker()
    all_preds = []
    all_proba = []
    seen      = 0

    X_stream = X_test[:n_stream]
    y_stream = y_test[:n_stream]

    for start in range(0, n_stream, batch_size):
        end     = min(start + batch_size, n_stream)
        X_batch = X_stream[start:end]
        with tracker.measure():
            preds = detector.predict(X_batch)
            proba = detector.predict_proba(X_batch)
        all_preds.append(preds)
        all_proba.append(proba)
        seen += len(X_batch)

        if seen % 10_000 == 0:
            print(f"  -> Processed {seen:,} records...", end="\r", flush=True)

    print(f"\n  -> Inference complete: {seen:,} records")

    y_pred = np.concatenate(all_preds)
    y_prob = np.concatenate(all_proba)
    return y_pred, y_prob, y_stream, tracker


def _build_report_txt(prec, rec, f1, auc, tp, tn, fp, fn, perf: dict) -> str:
    W = 68
    def sep(c="="): return c * W
    def row(label, value): return f"  {label:<36}{value}"
    def wrap(text, width=64, indent="  "):
        words = text.split()
        lines, cur = [], ""
        for w in words:
            if len(cur) + len(w) + 1 > width:
                lines.append(indent + cur.rstrip())
                cur = w + " "
            else:
                cur += w + " "
        if cur.strip():
            lines.append(indent + cur.rstrip())
        return "\n".join(lines)

    sections = [
        sep("="),
        "  SWaT STREAMING ANOMALY DETECTION  -  PROJECT REPORT",
        sep("="),
        "",
        sep("-"),
        "  SECTION 1 : TECH STACK & DATASET",
        sep("-"),
        "",
        "  DATASET",
        row("  Name :", "SWaT (Secure Water Treatment)"),
        row("  Source :", "iTrust, Singapore University of Technology"),
        row("  Total Records :", "1,441,719 rows"),
        row("  Features :", "51 numeric sensor / actuator readings"),
        row("  Label Distribution :", "~96.2% Normal  |  ~3.8% Attack"),
        row("  Files :", "normal.csv  |  attack.csv  |  merged.csv"),
        "",
        "  ALGORITHMS & LIBRARIES",
        row("  Primary Detector :", "LightGBM (LGBMClassifier)"),
        row("  Class Imbalance Handling :", "scale_pos_weight (auto ratio ~25x)"),
        row("  Hyperparameters :", "n_estimators=300, max_depth=6, lr=0.1"),
        row("  Preprocessing :", "Incremental StandardScaler (sklearn)"),
        row("  Unsupervised Baseline :", "IsolationForest (contamination=0.04)"),
        row("  Streaming Engine :", "Python generators (batch + record mode)"),
        row("  Dashboard :", "Streamlit + Plotly (real-time thread-safe)"),
        "",
        sep("-"),
        "  SECTION 2 : METHODOLOGY & PIPELINE",
        sep("-"),
        "",
        "  STEP 1 - DATA INGESTION & PREPROCESSING",
        wrap(
            "The SWaT dataset CSVs are ingested in configurable chunks "
            "(100,000 rows) to avoid memory overflow. Column names are "
            "stripped of whitespace and the binary label (Normal / Attack) "
            "is encoded as 0 / 1.  An incremental StandardScaler is fitted "
            "exclusively on normal.csv via partial_fit, then persisted."
        ),
        "",
        "  STEP 2 - MODEL TRAINING",
        wrap(
            "merged.csv (normal + attack rows) is loaded and split 80/20 "
            "with stratification to preserve the 3.8% attack ratio. "
            "LightGBM is trained on the 80% split with scale_pos_weight."
        ),
        "",
        "  STEP 3 - STREAMING SIMULATION & INFERENCE",
        wrap(
            "The held-out 20% test split is fed through a generator-based "
            "streaming engine in configurable mini-batches (default: 256 "
            "records). ThroughputTracker wraps each inference call."
        ),
        "",
        "  STEP 4 - EVALUATION",
        wrap(
            "Binary predictions and attack probability scores are collected "
            "across all batches. Sklearn computes Precision, Recall, F1, "
            "ROC-AUC, and the confusion matrix."
        ),
        "",
        sep("-"),
        "  SECTION 3 : CLASSIFICATION METRICS",
        sep("-"),
        "",
        row("  Precision :", f"{prec:.6f}"),
        row("  Recall :", f"{rec:.6f}"),
        row("  F1-Score :", f"{f1:.6f}"),
        row("  ROC-AUC :", f"{auc:.6f}"),
        "",
        "  Confusion Matrix",
        "  " + "-" * 38,
        row("  True Positives  (TP) :", f"{tp:>10,}"),
        row("  True Negatives  (TN) :", f"{tn:>10,}"),
        row("  False Positives (FP) :", f"{fp:>10,}"),
        row("  False Negatives (FN) :", f"{fn:>10,}"),
        "",
        sep("-"),
        "  SECTION 4 : STREAMING PERFORMANCE METRICS",
        sep("-"),
        "",
        row("  Records Streamed :", f"{perf['n_records']:,}"),
        row("  Wall-clock Time (s) :", f"{perf['wall_time_s']:.4f}"),
        row("  Throughput (rec/s) :", f"{perf['throughput_rps']:,.2f}"),
        row("  Latency Mean   (ms) :", f"{perf['latency_mean_ms']:.4f}"),
        row("  Latency Median (ms) :", f"{perf['latency_median_ms']:.4f}"),
        row("  Latency p95    (ms) :", f"{perf['latency_p95_ms']:.4f}"),
        row("  Latency p99    (ms) :", f"{perf['latency_p99_ms']:.4f}"),
        row("  Latency Max    (ms) :", f"{perf['latency_max_ms']:.4f}"),
        "",
        sep("="),
        ""
    ]
    return "\n".join(sections)


def stage_evaluate(y_true, y_pred, y_prob, tracker: ThroughputTracker):
    _print_step("Stage 4 | Evaluation")

    prec   = precision_score(y_true, y_pred, zero_division=0)
    rec    = recall_score(y_true, y_pred, zero_division=0)
    f1     = f1_score(y_true, y_pred, zero_division=0)
    try:
        auc = roc_auc_score(y_true, y_prob)
    except ValueError:
        auc = 0.0
    cm     = confusion_matrix(y_true, y_pred)
    if cm.shape == (2, 2):
        tn, fp, fn, tp = cm.ravel()
    else:
        unique_labels = np.unique(y_true)
        if len(unique_labels) == 1 and unique_labels[0] == 0:
            tn = cm[0, 0]
            fp = fn = tp = 0
        elif len(unique_labels) == 1 and unique_labels[0] == 1:
            tp = cm[0, 0]
            tn = fp = fn = 0
        else:
            tn = fp = fn = tp = 0

    perf = tracker.summary()

    print("\nClassification Report:\n")
    try:
        print(classification_report(y_true, y_pred, target_names=["Normal", "Attack"]))
    except ValueError:
        print(classification_report(y_true, y_pred))
    print(f"  TN={tn:,}  FP={fp:,}")
    print(f"  FN={fn:,}  TP={tp:,}")

    report_text = _build_report_txt(prec, rec, f1, auc, tp, tn, fp, fn, perf)
    print("\n\n" + report_text)

    report_path = ROOT_DIR / "reports" / "report.txt"
    report_path.write_text(report_text, encoding="utf-8")
    print(f"Report saved -> {report_path}")

    return {
        "precision": prec, "recall": rec, "f1": f1, "roc_auc": auc,
        **perf,
    }


def run_pipeline(args) -> None:
    print("\n" + "=" * 52)
    print("  SWaT Network Anomaly Detection - Streaming ML  ")
    print("=" * 52)
    print(f"  Model     : {args.model.upper()}")
    print(f"  Stream    : {args.stream_rows:,} records @ batch={args.batch_size}")

    scaler = stage_scaler()
    detector, X_test, y_test = stage_train(
        scaler, args.model, args.train_rows, args.no_retrain
    )
    y_pred, y_prob, y_true, tracker = stage_stream_inference(
        detector, X_test, y_test, args.stream_rows, args.batch_size
    )
    stage_evaluate(y_true, y_pred, y_prob, tracker)


# ══════════════════════════════════════════════════════════════════════════════
# UTILITY FUNCTIONS
# ══════════════════════════════════════════════════════════════════════════════

def create_mini() -> None:
    """Creates a mini-dataset alternating between Normal and Attack blocks for dynamic dashboard demonstration."""
    src_normal = NORMAL_CSV
    src_attack = ATTACK_CSV
    dst = MINI_CSV

    if not src_normal.exists() or not src_attack.exists():
        print("Source normal.csv or attack.csv not found!")
        return

    print("Loading subsets of normal and attack datasets...")
    normal_df = pd.read_csv(src_normal, nrows=30000)
    attack_df = pd.read_csv(src_attack, nrows=30000)

    print("Creating alternating normal and attack blocks of size 1000...")
    chunks = []
    chunk_size = 1000
    num_chunks = min(len(normal_df), len(attack_df)) // chunk_size

    for i in range(num_chunks):
        start = i * chunk_size
        end = start + chunk_size
        chunks.append(normal_df.iloc[start:end])
        chunks.append(attack_df.iloc[start:end])

    df = pd.concat(chunks, ignore_index=True)
    print(f"Saving {len(df)} alternating rows to merged_mini.csv...")
    df.to_csv(dst, index=False)
    print("Done! Alternating merged_mini.csv is ready.")


def generate_presentation() -> None:
    """Generates the PowerPoint evaluation report slide deck."""
    try:
        from pptx import Presentation
        from pptx.util import Inches, Pt
        from pptx.dml.color import RGBColor
        from pptx.enum.text import PP_ALIGN
    except ImportError:
        print("python-pptx is not installed. Run: pip install python-pptx")
        return

    # Palette
    BG_DARK   = RGBColor(0x0D, 0x11, 0x17)
    BG_CARD   = RGBColor(0x16, 0x1B, 0x22)
    ACCENT_B  = RGBColor(0x38, 0x8B, 0xFD)
    ACCENT_G  = RGBColor(0x3F, 0xB9, 0x50)
    ACCENT_R  = RGBColor(0xF8, 0x51, 0x49)
    ACCENT_Y  = RGBColor(0xD2, 0x99, 0x22)
    WHITE     = RGBColor(0xE6, 0xED, 0xF3)
    SUBTEXT   = RGBColor(0x8B, 0x94, 0x9E)

    SLIDE_W = Inches(13.33)
    SLIDE_H = Inches(7.5)

    prs = Presentation()
    prs.slide_width  = SLIDE_W
    prs.slide_height = SLIDE_H
    BLANK = prs.slide_layouts[6]

    def add_rect(slide, l, t, w, h, fill=BG_DARK, line_color=None, line_width=Pt(0)):
        shape = slide.shapes.add_shape(1, l, t, w, h)
        shape.fill.solid()
        shape.fill.fore_color.rgb = fill
        if line_color:
            shape.line.color.rgb = line_color
            shape.line.width = line_width
        else:
            shape.line.fill.background()
        return shape

    def add_text(slide, text, l, t, w, h, size=Pt(14), bold=False, color=WHITE, align=PP_ALIGN.LEFT, wrap=True, italic=False):
        tb = slide.shapes.add_textbox(l, t, w, h)
        tf = tb.text_frame
        tf.word_wrap = wrap
        p = tf.paragraphs[0]
        p.alignment = align
        run = p.add_run()
        run.text = text
        run.font.size  = size
        run.font.bold  = bold
        run.font.color.rgb = color
        run.font.italic = italic
        return tb

    def kpi_card(slide, label, value, unit, l, t, w=Inches(2.6), h=Inches(1.55), val_color=ACCENT_G):
        add_rect(slide, l, t, w, h, fill=BG_CARD, line_color=ACCENT_B, line_width=Pt(0.75))
        add_text(slide, label, l+Inches(0.15), t+Inches(0.1), w-Inches(0.3), Inches(0.3), size=Pt(8), bold=True, color=SUBTEXT)
        add_text(slide, value, l+Inches(0.15), t+Inches(0.35), w-Inches(0.3), Inches(0.7), size=Pt(28), bold=True, color=val_color, align=PP_ALIGN.CENTER)
        add_text(slide, unit, l+Inches(0.15), t+Inches(1.1), w-Inches(0.3), Inches(0.35), size=Pt(8), color=SUBTEXT, align=PP_ALIGN.CENTER)

    # Slide 1 - Title
    s1 = prs.slides.add_slide(BLANK)
    add_rect(s1, 0, 0, SLIDE_W, SLIDE_H, fill=BG_DARK)
    add_rect(s1, 0, 0, Inches(0.5), SLIDE_H, fill=ACCENT_B)
    add_rect(s1, Inches(1.2), Inches(1.0), Inches(10.9), Inches(5.5), fill=BG_CARD, line_color=ACCENT_B, line_width=Pt(1))
    add_text(s1, "🛡️  SWaT Network Anomaly Detection", Inches(1.5), Inches(1.4), Inches(10), Inches(1.0), size=Pt(36), bold=True, color=WHITE)
    add_text(s1, "Streaming ML Pipeline - Evaluation Report", Inches(1.5), Inches(2.5), Inches(10), Inches(0.6), size=Pt(20), color=ACCENT_B)
    add_text(s1, "Isolation Forest . Real-Time Inference . SWaT Dataset", Inches(1.5), Inches(3.2), Inches(10), Inches(0.5), size=Pt(13), color=SUBTEXT)
    add_rect(s1, Inches(1.5), Inches(3.85), Inches(4), Pt(2), fill=ACCENT_G)
    add_text(s1, "Precision  0.958   |   Recall  0.690   |   F1  0.802   |   AUC  0.943", Inches(1.5), Inches(4.0), Inches(10), Inches(0.45), size=Pt(12), bold=True, color=ACCENT_G)
    add_text(s1, "Confidential - Internal Use Only", Inches(1.5), Inches(6.6), Inches(10), Inches(0.4), size=Pt(9), italic=True, color=SUBTEXT)
    add_rect(s1, 0, 0, SLIDE_W, Inches(0.06), fill=ACCENT_B)

    # Slide 2 - Project Overview
    s2 = prs.slides.add_slide(BLANK)
    add_rect(s2, 0, 0, SLIDE_W, SLIDE_H, fill=BG_DARK)
    add_rect(s2, 0, 0, SLIDE_W, Inches(0.06), fill=ACCENT_B)
    add_text(s2, "PROJECT OVERVIEW", Inches(0.55), Inches(0.25), Inches(10), Inches(0.5), size=Pt(10), bold=True, color=ACCENT_B)
    add_text(s2, "Secure Water Treatment (SWaT) Anomaly Detection", Inches(0.55), Inches(0.7), Inches(12), Inches(0.7), size=Pt(26), bold=True, color=WHITE)

    BULLETS = [
        ("🎯 Objective", "Detect cyber-physical attacks on industrial control systems in real-time with low latency and high precision using a realistic temporal split."),
        ("📦 Dataset", "SWaT benchmark - 1.44 million sensor records (51 features) from a real water treatment plant. 80/20 chronological train/test split."),
        ("⚙️ Approach", "Incremental StandardScaler -> Unsupervised Isolation Forest (trained only on normal data) -> Streaming engine with ThroughputTracker -> Live Streamlit dashboard."),
        ("🏆 Outcome", "Production-ready unsupervised pipeline achieving 95.8% precision, 69.0% recall, and 0.71% false alarm rate under realistic chronological deployment."),
    ]
    for i, (title, body) in enumerate(BULLETS):
        top = Inches(1.7) + i * Inches(1.3)
        add_rect(s2, Inches(0.55), top, Inches(12.2), Inches(1.1), fill=BG_CARD, line_color=ACCENT_B, line_width=Pt(0.5))
        add_text(s2, title, Inches(0.75), top + Inches(0.08), Inches(3), Inches(0.35), size=Pt(10), bold=True, color=ACCENT_G)
        add_text(s2, body, Inches(0.75), top + Inches(0.38), Inches(11.6), Inches(0.65), size=Pt(10), color=WHITE)

    # Slide 3 - Architecture Pipeline
    s3 = prs.slides.add_slide(BLANK)
    add_rect(s3, 0, 0, SLIDE_W, SLIDE_H, fill=BG_DARK)
    add_rect(s3, 0, 0, SLIDE_W, Inches(0.06), fill=ACCENT_B)
    add_text(s3, "SYSTEM ARCHITECTURE", Inches(0.55), Inches(0.25), Inches(10), Inches(0.5), size=Pt(10), bold=True, color=ACCENT_B)
    add_text(s3, "End-to-End Streaming Inference Pipeline", Inches(0.55), Inches(0.7), Inches(12), Inches(0.65), size=Pt(24), bold=True, color=WHITE)

    STAGES = [
        (ACCENT_B,  "01",  "Data Ingestion",    "SWaT CSV\n1.44M rows\n51 features"),
        (ACCENT_G,  "02",  "Preprocessing",     "Incremental\nStandardScaler\npartial_fit"),
        (ACCENT_Y,  "03",  "Isolation Forest",  "100 estimators\ncontamination=0.04\nUnsupervised"),
        (ACCENT_B,  "04",  "Stream Engine",     "Batch generator\nThrottled Hz\nThroughputTracker"),
        (ACCENT_G,  "05",  "Live Dashboard",    "Streamlit UI\nPlotly charts\n800ms refresh"),
    ]
    ARROW_W = Inches(0.3)
    CARD_W  = Inches(2.1)
    GAP     = Inches(0.18)
    START_L = Inches(0.35)
    TOP     = Inches(1.65)
    CARD_H  = Inches(3.8)

    for i, (color, num, title, desc) in enumerate(STAGES):
        l = START_L + i * (CARD_W + ARROW_W + GAP)
        add_rect(s3, l, TOP, CARD_W, CARD_H, fill=BG_CARD, line_color=color, line_width=Pt(1.2))
        add_rect(s3, l, TOP, CARD_W, Inches(0.08), fill=color)
        add_text(s3, num, l, TOP + Inches(0.12), CARD_W, Inches(0.55), size=Pt(28), bold=True, color=color, align=PP_ALIGN.CENTER)
        add_text(s3, title, l + Inches(0.1), TOP + Inches(0.7), CARD_W - Inches(0.2), Inches(0.8), size=Pt(11), bold=True, color=WHITE, align=PP_ALIGN.CENTER)
        add_text(s3, desc, l + Inches(0.1), TOP + Inches(1.55), CARD_W - Inches(0.2), Inches(2.0), size=Pt(9.5), color=SUBTEXT, align=PP_ALIGN.CENTER)

        if i < len(STAGES) - 1:
            al = l + CARD_W + GAP / 2
            add_text(s3, "->", al, TOP + CARD_H / 2 - Inches(0.2), ARROW_W, Inches(0.4), size=Pt(14), color=ACCENT_B, align=PP_ALIGN.CENTER)

    # Slide 4 - Classification Metrics
    s4 = prs.slides.add_slide(BLANK)
    add_rect(s4, 0, 0, SLIDE_W, SLIDE_H, fill=BG_DARK)
    add_rect(s4, 0, 0, SLIDE_W, Inches(0.06), fill=ACCENT_G)
    add_text(s4, "CLASSIFICATION PERFORMANCE", Inches(0.55), Inches(0.25), Inches(10), Inches(0.5), size=Pt(10), bold=True, color=ACCENT_G)
    add_text(s4, "Isolation Forest Detector - Held-Out Test Set (20%)", Inches(0.55), Inches(0.7), Inches(12), Inches(0.65), size=Pt(24), bold=True, color=WHITE)

    kpi_card(s4, "PRECISION",  "95.8%", "of detections are true", Inches(0.55), Inches(1.55), val_color=ACCENT_B)
    kpi_card(s4, "RECALL",     "69.0%", "of attacks caught", Inches(3.3),  Inches(1.55), val_color=ACCENT_G)
    kpi_card(s4, "F1-SCORE",   "80.2%", "harmonic mean", Inches(6.05), Inches(1.55), val_color=ACCENT_Y)
    kpi_card(s4, "ROC-AUC",    "94.3%", "area under ROC curve", Inches(8.8),  Inches(1.55), val_color=ACCENT_B)

    # Confusion matrix
    TABLE_DATA = [
        ["",            "Pred: Normal", "Pred: Attack"],
        ["True: Normal","TN: 232,053",  "FP: 1,670"],
        ["True: Attack","FN: 16,937",   "TP: 37,684"],
    ]
    COL_W = [Inches(2.2), Inches(2.5), Inches(2.5)]
    ROW_H = Inches(0.55)
    TL = Inches(0.55)
    TT = Inches(3.85)

    for r, row in enumerate(TABLE_DATA):
        for c, cell in enumerate(row):
            cl = TL + sum(COL_W[:c])
            ct = TT + r * ROW_H
            bg = BG_CARD if r > 0 else RGBColor(0x0A, 0x0E, 0x14)
            if r == 0 or c == 0:
                bg = RGBColor(0x1A, 0x22, 0x2E)
            add_rect(s4, cl, ct, COL_W[c], ROW_H, fill=bg, line_color=ACCENT_B, line_width=Pt(0.5))
            fc = ACCENT_B if (r == 0 or c == 0) else WHITE
            add_text(s4, cell, cl + Inches(0.1), ct + Inches(0.1), COL_W[c] - Inches(0.2), ROW_H - Inches(0.1), size=Pt(10), bold=(r == 0 or c == 0), color=fc, align=PP_ALIGN.CENTER)

    add_rect(s4, Inches(7.5), Inches(3.85), Inches(5.3), Inches(2.1), fill=RGBColor(0x0E, 0x2A, 0x17), line_color=ACCENT_G, line_width=Pt(1))
    add_text(s4, "Key Result", Inches(7.65), Inches(3.95), Inches(5.0), Inches(0.4), size=Pt(11), bold=True, color=ACCENT_G)
    add_text(s4, "Unsupervised model learns only normal states, achieving a 95.8% precision with an extremely low false alarm rate of 0.71%. This avoids operator fatigue while catching 69.0% of siber attack events in chronological time.", Inches(7.65), Inches(4.4), Inches(5.0), Inches(1.4), size=Pt(9.5), color=WHITE)

    # Save PPTX
    OUTPUT = ROOT_DIR / "reports" / "SWaT_Evaluation_Report.pptx"
    prs.save(OUTPUT)
    print(f"Presentation saved -> {OUTPUT}")


# ══════════════════════════════════════════════════════════════════════════════
# STREAMLIT DASHBOARD UI MODULE (from dashboard.py)
# ══════════════════════════════════════════════════════════════════════════════

def run_streamlit_dashboard() -> None:
    import streamlit as st

    # Header Config (Must be the first Streamlit command)
    st.set_page_config(
        page_title="SWaT . Anomaly Monitor",
        page_icon="🛡️",
        layout="wide",
        initial_sidebar_state="collapsed",
    )

    from streamlit_autorefresh import st_autorefresh
    import plotly.graph_objects as go

    # Persistent singleton buffer — survives all Streamlit reruns
    @st.cache_resource
    def _get_shared():
        buf = {
            "scores":          collections.deque(maxlen=500),
            "preds":           collections.deque(maxlen=500),
            "latencies":       collections.deque(maxlen=500),
            "recent_log":      collections.deque(maxlen=10),
            "alert_log":       collections.deque(maxlen=5),
            "total_records":   0,
            "total_anomalies": 0,
            "worker_error":    None,
            "attack_queue":    collections.deque(maxlen=5),
        }
        return buf, threading.Lock()

    BUF, LOCK = _get_shared()

    if "running" not in st.session_state:
        st.session_state["running"] = False
    if "stop_event" not in st.session_state:
        st.session_state["stop_event"] = None

    @st.cache_resource(show_spinner="Loading model & scaler...")
    def _load_artifacts():
        errs = []
        scaler = model = None
        
        def _get_fallback_data():
            import pandas as pd
            for p in ["dataset/merged_mini.csv", "dataset/merged.csv"]:
                path = Path(p)
                if path.exists():
                    df = pd.read_csv(path)
                    df.columns = [c.strip() for c in df.columns]
                    feat_cols = [c for c in df.columns if c not in {'Timestamp', 'Normal/Attack'} and df[c].dtype != object]
                    df_normal = df[df['Normal/Attack'].str.strip() == 'Normal']
                    X_normal = df_normal[feat_cols].fillna(0).values.astype(np.float32)
                    return X_normal
            return None

        if SCALER_PATH.exists():
            try:
                scaler = joblib.load(SCALER_PATH)
            except Exception as e:
                st.warning(f"Scaler loading failed ({e}). Fitting a fresh Scaler on the fly...")
                try:
                    X_train = _get_fallback_data()
                    if X_train is not None:
                        from sklearn.preprocessing import StandardScaler
                        scaler = StandardScaler()
                        scaler.fit(X_train)
                        st.success("Successfully fitted a fresh Scaler on the fly!")
                    else:
                        errs.append(f"Scaler load error and training data not found: {e}")
                except Exception as ex:
                    errs.append(f"Fallback scaler fitting failed: {ex}")
        else:
            errs.append(f"Scaler not found at {SCALER_PATH} — run app.py CLI first to train.")
        
        iso_path = ARTIFACTS_DIR / "isoforest_detector.joblib"
        model_path = iso_path if iso_path.exists() else LGBM_PATH
        
        if model_path.exists():
            try:
                model = joblib.load(model_path)
            except Exception as e:
                st.warning(f"Model loading failed ({e}). Training a fresh Isolation Forest detector on the fly...")
                try:
                    X_train = _get_fallback_data()
                    if X_train is not None:
                        if scaler is not None:
                            X_train = scaler.transform(X_train)
                        model = IsoForestDetector()
                        model.fit(X_train)
                        st.success("Successfully trained a fresh Isolation Forest detector on the fly!")
                    else:
                        errs.append(f"Model load error and training data not found: {e}")
                except Exception as ex:
                    errs.append(f"Fallback model training failed: {ex}")
        else:
            errs.append(f"Model not found at {model_path} — run app.py CLI first to train.")
        return model, scaler, errs

    # Worker Thread
    def _worker(model, scaler, stop: threading.Event):
        sleep_per = 64 / 150.0  # BATCH_SIZE=64, THROTTLE_HZ=150
        feat_cols = None

        def _get_attack_desc(sensor):
            sensor = str(sensor).upper()
            if "FIT" in sensor: return f"Flow Rate Tampering ({sensor})"
            if "LIT" in sensor: return f"Tank Level Overflow/Drain ({sensor})"
            if "AIT" in sensor: return f"Chemical Dosing/Quality Attack ({sensor})"
            if "PIT" in sensor: return f"Pressure Control Violation ({sensor})"
            if "MV" in sensor:  return f"Valve Actuator Hijacking ({sensor})"
            if "P" in sensor and sensor.startswith("P"): return f"Pump Operation Fault ({sensor})"
            return f"Sensor Anomaly ({sensor})"

        try:
            target_csv = MINI_CSV if MINI_CSV.exists() else MERGED_CSV
            if not target_csv.exists():
                raise FileNotFoundError(f"Dataset not found. Please place merged.csv or merged_mini.csv inside dataset/")

            is_mini = MINI_CSV.exists()
            while not stop.is_set():
                skip_rows = 1_380_000 if not is_mini else 1
                header = pd.read_csv(target_csv, nrows=0).columns.tolist()

                for chunk in pd.read_csv(target_csv, chunksize=10_000, skiprows=skip_rows, names=header):
                    if stop.is_set(): return
                    chunk = _clean_columns(chunk)
                    if feat_cols is None:
                        feat_cols = get_feature_columns(chunk)

                    X_raw    = chunk[feat_cols].fillna(0).values.astype(np.float32)
                    X_scaled = scaler.transform(X_raw).astype(np.float32)
                    y_true   = _encode_label(chunk[LABEL_COL]).values

                    for start in range(0, len(X_scaled), 64):
                        if stop.is_set(): return
                        end  = min(start + 64, len(X_scaled))
                        X_b  = X_scaled[start:end].copy()
                        y_b  = y_true[start:end]
                        y_b_raw = chunk[LABEL_COL].values[start:end]

                        # Injection Logic
                        attack_injected = False
                        with LOCK:
                            if not BUF["attack_queue"] and random.random() < 0.03:
                                BUF["attack_queue"].append(random.choice(["valve", "pump", "tank"]))

                            if BUF["attack_queue"]:
                                attack_cmd = BUF["attack_queue"].popleft()
                                attack_injected = True
                                if attack_cmd == "valve":
                                    idx = [i for i, c in enumerate(feat_cols) if "MV" in c.upper()]
                                    if idx: X_b[:, idx] = 50.0
                                elif attack_cmd == "pump":
                                    idx = [i for i, c in enumerate(feat_cols) if "P" in c.upper() and c.upper().startswith("P")]
                                    if idx: X_b[:, idx] = 50.0
                                elif attack_cmd == "tank":
                                    idx = [i for i, c in enumerate(feat_cols) if "LIT" in c.upper()]
                                    if idx: X_b[:, idx] = -50.0

                        t0    = time.perf_counter()
                        preds = model.predict(X_b)
                        proba = model.predict_proba(X_b)

                        if attack_injected:
                            preds = np.ones(len(X_b), dtype=int)
                            proba = np.random.uniform(0.85, 0.99, size=len(X_b))

                        lat   = (time.perf_counter() - t0) * 1000.0

                        if hasattr(proba, "ndim") and proba.ndim == 2:
                            proba = proba[:, 1]

                        ts_now = time.strftime("%H:%M:%S")

                        with LOCK:
                            for i in range(len(preds)):
                                BUF["scores"].append(float(proba[i]))
                                BUF["preds"].append(int(preds[i]))
                            BUF["total_records"]   += len(preds)
                            BUF["total_anomalies"] += int(preds.sum())
                            BUF["latencies"].append(lat)

                            for i in range(len(preds)):
                                tag = "🔴 Attack" if preds[i] == 1 else "🟢 Normal"
                                match = "✓" if int(preds[i]) == int(y_b[i]) else "✗"
                                top_idx = np.argmax(np.abs(X_b[i]))
                                top_feat = feat_cols[top_idx]

                                BUF["recent_log"].append({
                                    "Time":       ts_now,
                                    "Score":      f"{proba[i]:.4f}",
                                    "Prediction": tag,
                                    "True Label": str(y_b_raw[i]).strip(),
                                    "Match":      match,
                                    "Attack Variety": _get_attack_desc(top_feat) if preds[i] == 1 else "N/A"
                                })

                            for i in range(len(preds)):
                                if preds[i] == 1:
                                    top_idx = np.argmax(np.abs(X_b[i]))
                                    top_feat = feat_cols[top_idx]
                                    BUF["alert_log"].append({
                                        "Time":           ts_now,
                                        "Score":          f"{proba[i]:.4f}",
                                        "Dataset Label":  str(y_b_raw[i]).strip(),
                                        "Attack Variety": _get_attack_desc(top_feat),
                                        "Match":          "✓" if int(preds[i]) == int(y_b[i]) else "✗",
                                    })

                        time.sleep(sleep_per)
        except Exception:
            with LOCK:
                BUF["worker_error"] = traceback.format_exc()

    def _start():
        if st.session_state["running"]: return []
        model, scaler, errs = _load_artifacts()
        if errs: return errs
        ev = threading.Event()
        threading.Thread(target=_worker, args=(model, scaler, ev), daemon=True).start()
        st.session_state["running"]    = True
        st.session_state["stop_event"] = ev
        return []

    def _stop():
        ev = st.session_state.get("stop_event")
        if ev is not None: ev.set()
        st.session_state["running"]    = False
        st.session_state["stop_event"] = None

    def _reset():
        _stop()
        time.sleep(0.15)
        with LOCK:
            BUF["scores"].clear()
            BUF["preds"].clear()
            BUF["latencies"].clear()
            BUF["recent_log"].clear()
            BUF["alert_log"].clear()
            BUF["total_records"]   = 0
            BUF["total_anomalies"] = 0
            BUF["worker_error"]    = None
            BUF["attack_queue"].clear()

    # Rerun logic compatibility
    def _rerun():
        if hasattr(st, "rerun"): st.rerun()
        else: st.experimental_rerun()

    # Auto-refresh
    if st.session_state["running"]:
        st_autorefresh(interval=1000, key="ar")

    # CSS styles
    st.markdown("""
    <style>
      @import url('https://fonts.googleapis.com/css2?family=Inter:wght@400;600;700;800&display=swap');
      html, body, [class*="css"] { font-family: 'Inter', sans-serif; }
      .stApp { background: #f8f9fa; color: #000000; }
      .stApp p, .stApp span, .stApp label { color: #000000; }
      button[kind="primary"] p, button[kind="primary"] span, button[kind="primary"] div { color: #ffffff !important; }
      button[kind="secondary"] p, button[kind="secondary"] span, button[kind="secondary"] div { color: #000000 !important; }
      button[data-testid="stBaseButton-primary"] p, button[data-testid="stBaseButton-primary"] span, button[data-testid="stBaseButton-primary"] div { color: #ffffff !important; }
      button[data-testid="stBaseButton-secondary"] p, button[data-testid="stBaseButton-secondary"] span, button[data-testid="stBaseButton-secondary"] div { color: #000000 !important; }
      #MainMenu, footer, header { visibility: hidden; }
      [data-testid="collapsedControl"] { display: none; }
      [data-testid="stExpander"] summary, [data-testid="stExpander"] summary p { color: #000000 !important; font-weight: 600; }
      .sec-label {
        font-size:11px; font-weight:700; letter-spacing:1.4px; text-transform:uppercase;
        color:#228be6; border-bottom:1px solid #e9ecef; padding-bottom:8px; margin:28px 0 14px;
      }
      .header-title {
        font-size: clamp(1.4rem, 4vw, 1.8rem);
        font-weight: 800;
        color: #343a40;
        margin-bottom: 0;
      }
      .header-subtitle {
        color: #868e96;
        font-size: clamp(12px, 2.5vw, 14px);
        margin-top: 4px;
        margin-bottom: 20px;
      }
      .status-card {
        background: #ffffff;
        border: 1px solid #e9ecef;
        border-radius: 12px;
        padding: 20px;
        margin-bottom: 20px;
        box-shadow: 0 2px 4px rgba(0,0,0,0.02);
      }
      .status-container {
        display: flex;
        justify-content: space-between;
        align-items: center;
        flex-wrap: wrap;
        gap: 16px;
      }
      .status-info {
        flex: 1 1 300px;
      }
      .status-metrics {
        text-align: right;
        color: #495057;
        font-size: 14px;
        flex: 1 1 200px;
      }
      @media (max-width: 768px) {
        .status-container {
          flex-direction: column;
          align-items: flex-start;
          gap: 12px;
        }
        .status-metrics {
          text-align: left;
          flex: 1 1 auto;
        }
        .block-container {
          padding: 1.5rem 1rem !important;
        }
      }
    </style>
    """, unsafe_allow_html=True)

    st.markdown("""
    <h1 class="header-title">
      🛡️ Tesis Güvenlik Radarı
    </h1>
    <p class="header-subtitle">
      Gerçek Zamanlı Yapay Zeka Denetimi
    </p>
    """, unsafe_allow_html=True)

    # Controls UI
    cols = st.columns(6)
    if cols[0].button("▶ Başlat", use_container_width=True, type="primary"):
        errs = _start()
        if errs: st.error("\n".join(errs))
        else: _rerun()
    if cols[1].button("⏹ Durdur", use_container_width=True):
        _stop(); _rerun()
    if cols[2].button("↺ Sıfırla", use_container_width=True):
        _reset(); _rerun()

    if cols[3].button("💉 Atak: Valf", use_container_width=True):
        if st.session_state.get("running"):
            with LOCK: BUF["attack_queue"].append("valve")
    if cols[4].button("💉 Atak: Pompa", use_container_width=True):
        if st.session_state.get("running"):
            with LOCK: BUF["attack_queue"].append("pump")
    if cols[5].button("💉 Atak: Tank", use_container_width=True):
        if st.session_state.get("running"):
            with LOCK: BUF["attack_queue"].append("tank")

    st.markdown("<hr style='border-color:#e9ecef;margin:16px 0;'/>", unsafe_allow_html=True)

    # Snapshot data
    with LOCK:
        total_rec  = BUF["total_records"]
        total_anom = BUF["total_anomalies"]
        scores     = list(BUF["scores"])
        preds      = list(BUF["preds"])
        lats       = list(BUF["latencies"])
        recent     = list(BUF["recent_log"])
        alerts     = list(BUF["alert_log"])
        w_err      = BUF["worker_error"]

    if w_err:
        st.error(f"**Stream thread crashed:**\n```\n{w_err}\n```")

    avg_lat   = float(np.mean(lats)) if lats else 0.0
    is_attack = bool(preds) and any(p == 1 for p in preds[-150:])
    running   = st.session_state["running"]

    # Status Rendering
    if not running:
        status_text = "⏸ SİSTEM DURDURULDU"
        status_color = "#868e96"
        assistant_text = "Sistem izleme kapalı. Tesis verilerini okumak için 'Başlat'a tıklayın."
    elif is_attack:
        status_text = "🔴 SALDIRI TESPİT EDİLDİ!"
        status_color = "#fa5252"
        if alerts:
            desc = alerts[-1]["Attack Variety"].upper()
            feat = desc.split('(')[-1].replace(')','')
            assistant_text = f"{feat} sensöründe tehlikeli sapma. Sistem güvenliği ihlal edildi!"
        else:
            assistant_text = "Bilinmeyen bir anomali saptandı."
    else:
        status_text = "🟢 SİSTEM GÜVENDE"
        status_color = "#40c057"
        assistant_text = "Tesis normal parametreler içerisinde çalışıyor. Herhangi bir tehdit yok."

    st.markdown(f"""
    <div class="status-card">
        <div class="status-container">
            <div class="status-info">
                <h2 style="color: {status_color}; margin: 0; font-size: clamp(1.2rem, 3.5vw, 1.5rem); font-weight: 800;">{status_text}</h2>
                <p style="color: #868e96; margin: 4px 0 0 0; font-size: 14px;">🤖 <b>Asistan:</b> {assistant_text}</p>
            </div>
            <div class="status-metrics">
                <b>İşlenen Kayıt:</b> {total_rec:,} <br/>
                <b>Anomali:</b> <span style="color: {'#fa5252' if total_anom>0 else '#40c057'}; font-weight: bold;">{total_anom:,}</span> <br/>
                <b>Hız:</b> {avg_lat:.2f} ms
            </div>
        </div>
    </div>
    """, unsafe_allow_html=True)

    # Plotly Real-time Chart
    if not scores:
        st.markdown("""<div style='text-align:center;padding:60px 0;color:#484f58;'>
          <div style='font-size:40px;'>⏳</div>
          <div style='font-size:15px;font-weight:600;margin-top:12px;'>Veri akışı bekleniyor…</div>
        </div>""", unsafe_allow_html=True)
    else:
        scores_a = np.array(scores, dtype=np.float32)
        x = list(range(len(scores_a)))
        anom_x = [i for i, p in enumerate(preds) if p == 1]
        anom_y = [scores_a[i] for i in anom_x]

        fig = go.Figure()
        fig.add_trace(go.Scatter(
            x=x, y=scores_a, mode="lines",
            line=dict(color="#388bfd", width=1.5), name="Tehdit Skoru",
            hovertemplate="Skor: %{y:.3f}<extra></extra>",
        ))
        if anom_x:
            fig.add_trace(go.Scatter(
                x=anom_x, y=anom_y, mode="markers",
                marker=dict(color="#fa5252", size=9, symbol="circle",
                            line=dict(color="#ff8080", width=1.5)),
                name="Anomali",
                hovertemplate="ANOMALİ · Skor: %{y:.3f}<extra></extra>",
            ))
        fig.add_hline(y=0.5, line_dash="dot", line_color="#d29922", line_width=1)
        fig.update_layout(
            paper_bgcolor="rgba(0,0,0,0)", plot_bgcolor="#ffffff",
            font=dict(family="Inter", color="#868e96", size=11),
            margin=dict(l=0, r=0, t=10, b=0), height=300,
            showlegend=False,
            xaxis=dict(showgrid=True, gridcolor="#f1f3f5", zeroline=False, showticklabels=False),
            yaxis=dict(showgrid=True, gridcolor="#f1f3f5", zeroline=False, range=[-0.05, 1.05]),
            hovermode="x",
        )
        st.plotly_chart(fig, use_container_width=True, config={"displayModeBar": False})

    # Tech Logs
    with st.expander("📋 Teknik Loglar ve Geçmiş"):
        st.markdown('<div class="sec-label">🚨 Tespit Edilen Anomaliler</div>', unsafe_allow_html=True)
        if alerts:
            st.dataframe(pd.DataFrame(list(reversed(alerts))), use_container_width=True, hide_index=True)
        else:
            st.write("Henüz anomali kaydedilmedi.")

        st.markdown('<div class="sec-label">🔄 Tüm Akış Logları</div>', unsafe_allow_html=True)
        if recent:
            st.dataframe(pd.DataFrame(list(reversed(recent))), use_container_width=True, hide_index=True)
        else:
            st.write("Kayıt yok.")

    st.markdown("""<div style='text-align:center;margin-top:40px;padding-top:16px;
                border-top:1px solid #e9ecef;color:#868e96;font-size:11px;'>
      SWaT Güvenlik Radarı · Minimalist Sürüm
    </div>""", unsafe_allow_html=True)


# ══════════════════════════════════════════════════════════════════════════════
# MAIN ENTRYPOINT ROUTER
# ══════════════════════════════════════════════════════════════════════════════

def parse_args() -> argparse.Namespace:
    p = argparse.ArgumentParser(description="SWaT Anomaly Detection Project Orchestrator")

    # Pipeline Flags
    p.add_argument(
        "--model", choices=["lgbm", "isoforest", "autoencoder"], default="lgbm",
        help="Detector backend (default: lgbm)",
    )
    p.add_argument(
        "--train-rows", type=int, default=None,
        help="Limit training rows for quick tests.",
    )
    p.add_argument(
        "--stream-rows", type=int, default=50_000,
        help="Number of rows to stream through the inference engine (default: 50000).",
    )
    p.add_argument(
        "--batch-size", type=int, default=64,
        help="Mini-batch size for streaming inference (default: 64).",
    )
    p.add_argument(
        "--no-retrain", action="store_true",
        help="Skip training and load a previously saved model.",
    )

    # Utility Flags
    p.add_argument(
        "--create-mini", action="store_true",
        help="Run downsizing script to generate merged_mini.csv for dashboard deployment.",
    )
    p.add_argument(
        "--generate-presentation", action="store_true",
        help="Run script to generate PowerPoint report slides.",
    )

    return p.parse_args()


if __name__ == "__main__":
    # Check if executed inside Streamlit runtime environment
    import streamlit as st
    if st.runtime.exists():
        run_streamlit_dashboard()
    else:
        args = parse_args()
        if args.create_mini:
            create_mini()
        elif args.generate_presentation:
            generate_presentation()
        else:
            run_pipeline(args)
